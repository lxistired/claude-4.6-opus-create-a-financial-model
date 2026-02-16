"""Tests for markitwrite - AI virtual clipboard."""

from __future__ import annotations

import base64
import io
import os
import tempfile

import pytest
from PIL import Image

from markitwrite import MarkItWrite, WriteResult
from markitwrite._base_writer import DocumentWriter
from markitwrite.writers._docx_writer import DocxWriter
from markitwrite.writers._pptx_writer import PptxWriter
from markitwrite.writers._markdown_writer import MarkdownWriter


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_test_image(width: int = 200, height: int = 150, color: str = "red") -> bytes:
    """Create a minimal PNG image in memory."""
    img = Image.new("RGB", (width, height), color=color)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _make_test_image_stream(**kwargs) -> io.BytesIO:
    data = _make_test_image(**kwargs)
    stream = io.BytesIO(data)
    stream.seek(0)
    return stream


# ---------------------------------------------------------------------------
# WriteResult tests
# ---------------------------------------------------------------------------

class TestWriteResult:
    def test_to_stream(self):
        result = WriteResult(output=b"hello", target_format=".txt")
        stream = result.to_stream()
        assert stream.read() == b"hello"

    def test_save(self, tmp_path):
        result = WriteResult(output=b"data", target_format=".bin")
        path = str(tmp_path / "out.bin")
        result.save(path)
        with open(path, "rb") as f:
            assert f.read() == b"data"


# ---------------------------------------------------------------------------
# DocumentWriter base class tests
# ---------------------------------------------------------------------------

class TestDocumentWriter:
    def test_accepts_returns_false_for_unknown(self):
        writer = DocumentWriter()
        assert writer.accepts(".xyz") is False

    def test_insert_image_raises(self):
        writer = DocumentWriter()
        with pytest.raises(NotImplementedError):
            writer.insert_image(io.BytesIO(b""))


# ---------------------------------------------------------------------------
# DocxWriter tests
# ---------------------------------------------------------------------------

class TestDocxWriter:
    def test_accepts(self):
        w = DocxWriter()
        assert w.accepts(".docx") is True
        assert w.accepts("docx") is True
        assert w.accepts(".pptx") is False

    def test_create_new_docx(self):
        w = DocxWriter()
        result = w.insert_image(
            image_stream=_make_test_image_stream(),
            size={"width": 4.0},
        )
        assert isinstance(result, WriteResult)
        assert result.target_format == ".docx"
        assert len(result.output) > 0
        # Verify it's a valid DOCX (ZIP with specific signature)
        assert result.output[:2] == b"PK"

    def test_insert_into_existing_docx(self):
        from docx import Document

        # Create a document with some text
        doc = Document()
        doc.add_paragraph("Hello World")
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)

        w = DocxWriter()
        result = w.insert_image(
            image_stream=_make_test_image_stream(),
            document_stream=buf,
        )
        assert result.output[:2] == b"PK"
        assert len(result.output) > len(buf.getvalue())

    def test_insert_at_paragraph(self):
        from docx import Document

        doc = Document()
        doc.add_paragraph("Paragraph 0")
        doc.add_paragraph("Paragraph 1")
        doc.add_paragraph("Paragraph 2")
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)

        w = DocxWriter()
        result = w.insert_image(
            image_stream=_make_test_image_stream(),
            document_stream=buf,
            position={"paragraph": 1},
        )
        # Verify the result is valid
        assert result.output[:2] == b"PK"

    def test_save_to_file(self, tmp_path):
        writer = MarkItWrite()
        img_path = str(tmp_path / "test.png")
        with open(img_path, "wb") as f:
            f.write(_make_test_image())

        target = str(tmp_path / "output.docx")
        result = writer.paste(img_path, target=target)
        assert os.path.isfile(target)
        assert os.path.getsize(target) > 0


# ---------------------------------------------------------------------------
# PptxWriter tests
# ---------------------------------------------------------------------------

class TestPptxWriter:
    def test_accepts(self):
        w = PptxWriter()
        assert w.accepts(".pptx") is True
        assert w.accepts(".docx") is False

    def test_create_new_pptx(self):
        w = PptxWriter()
        result = w.insert_image(
            image_stream=_make_test_image_stream(),
            size={"width": 6.0},
        )
        assert isinstance(result, WriteResult)
        assert result.target_format == ".pptx"
        assert result.output[:2] == b"PK"

    def test_insert_into_existing_pptx(self):
        from pptx import Presentation

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        prs.slides.add_slide(slide_layout)
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        w = PptxWriter()
        result = w.insert_image(
            image_stream=_make_test_image_stream(),
            document_stream=buf,
        )
        assert result.output[:2] == b"PK"

    def test_insert_on_specific_slide(self):
        from pptx import Presentation

        prs = Presentation()
        for _ in range(3):
            prs.slides.add_slide(prs.slide_layouts[0])
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        w = PptxWriter()
        result = w.insert_image(
            image_stream=_make_test_image_stream(),
            document_stream=buf,
            position={"slide": 2},
        )
        # Verify it's valid and slide count didn't change (inserted on existing)
        result_prs = Presentation(io.BytesIO(result.output))
        assert len(result_prs.slides) == 3

    def test_save_to_file(self, tmp_path):
        writer = MarkItWrite()
        img_path = str(tmp_path / "test.png")
        with open(img_path, "wb") as f:
            f.write(_make_test_image())

        target = str(tmp_path / "output.pptx")
        result = writer.paste(img_path, target=target)
        assert os.path.isfile(target)


# ---------------------------------------------------------------------------
# MarkdownWriter tests
# ---------------------------------------------------------------------------

class TestMarkdownWriter:
    def test_accepts(self):
        w = MarkdownWriter()
        assert w.accepts(".md") is True
        assert w.accepts(".markdown") is True
        assert w.accepts(".docx") is False

    def test_create_new_markdown_embedded(self):
        w = MarkdownWriter()
        img_data = _make_test_image()
        result = w.insert_image(
            image_stream=io.BytesIO(img_data),
            embed=True,
            mime_type="image/png",
        )
        content = result.output.decode("utf-8")
        assert "![image](data:image/png;base64," in content
        # Verify the base64 is valid
        b64_part = content.split("base64,")[1].rstrip(")\n")
        decoded = base64.b64decode(b64_part)
        assert decoded == img_data

    def test_create_new_markdown_reference(self):
        w = MarkdownWriter()
        result = w.insert_image(
            image_stream=_make_test_image_stream(),
            embed=False,
            image_filename="chart.png",
        )
        content = result.output.decode("utf-8")
        assert "![image](images/chart.png)" in content

    def test_append_to_existing_markdown(self):
        existing = b"# Title\n\nSome text here.\n"
        w = MarkdownWriter()
        result = w.insert_image(
            image_stream=_make_test_image_stream(),
            document_stream=io.BytesIO(existing),
            embed=True,
        )
        content = result.output.decode("utf-8")
        assert content.startswith("# Title")
        assert "Some text here." in content
        assert "![image](data:image/png;base64," in content

    def test_insert_after_heading(self):
        existing = b"# Intro\n\nText.\n\n# Charts\n\nMore text.\n"
        w = MarkdownWriter()
        result = w.insert_image(
            image_stream=_make_test_image_stream(),
            document_stream=io.BytesIO(existing),
            position={"after_heading": "Charts"},
            embed=True,
        )
        content = result.output.decode("utf-8")
        lines = content.split("\n")
        # Find the Charts heading and verify image is right after it
        charts_idx = None
        for i, line in enumerate(lines):
            if "# Charts" in line:
                charts_idx = i
                break
        assert charts_idx is not None
        # Image should be within the next few lines
        nearby = "\n".join(lines[charts_idx : charts_idx + 4])
        assert "data:image/png;base64," in nearby

    def test_save_to_file(self, tmp_path):
        writer = MarkItWrite()
        img_path = str(tmp_path / "test.png")
        with open(img_path, "wb") as f:
            f.write(_make_test_image())

        target = str(tmp_path / "output.md")
        result = writer.paste(img_path, target=target, embed=True)
        assert os.path.isfile(target)
        with open(target, "r") as f:
            content = f.read()
        assert "![image](data:image/png;base64," in content


# ---------------------------------------------------------------------------
# MarkItWrite dispatcher tests
# ---------------------------------------------------------------------------

class TestMarkItWrite:
    def test_supported_formats(self):
        writer = MarkItWrite()
        formats = writer.supported_formats()
        assert ".docx" in formats
        assert ".pptx" in formats
        assert ".md" in formats

    def test_paste_from_path(self, tmp_path):
        img_path = str(tmp_path / "img.png")
        with open(img_path, "wb") as f:
            f.write(_make_test_image())

        writer = MarkItWrite()
        result = writer.paste(img_path, target_format=".md", embed=True)
        assert result.target_format == ".md"

    def test_paste_from_bytes(self):
        writer = MarkItWrite()
        result = writer.paste(
            _make_test_image(),
            target_format=".docx",
        )
        assert result.target_format == ".docx"
        assert result.output[:2] == b"PK"

    def test_paste_from_stream(self):
        writer = MarkItWrite()
        result = writer.paste(
            _make_test_image_stream(),
            target_format=".pptx",
        )
        assert result.target_format == ".pptx"

    def test_format_inferred_from_target(self, tmp_path):
        img_path = str(tmp_path / "img.png")
        with open(img_path, "wb") as f:
            f.write(_make_test_image())

        writer = MarkItWrite()
        target = str(tmp_path / "out.pptx")
        result = writer.paste(img_path, target=target)
        assert result.target_format == ".pptx"
        assert os.path.isfile(target)

    def test_unsupported_format_raises(self):
        writer = MarkItWrite()
        with pytest.raises(ValueError, match="No writer found"):
            writer.paste(_make_test_image(), target_format=".xyz")

    def test_roundtrip_with_real_screenshot(self):
        """Test using the actual screenshot files from the repo if available."""
        screenshot_dir = "/home/user/claude-opus-4.6-create-a-financial-model/screenshots"
        screenshot = os.path.join(screenshot_dir, "02-key-summary.png")
        if not os.path.isfile(screenshot):
            pytest.skip("Screenshot not available")

        writer = MarkItWrite()
        with tempfile.TemporaryDirectory() as tmp:
            # Paste into DOCX
            docx_path = os.path.join(tmp, "report.docx")
            result = writer.paste(screenshot, target=docx_path)
            assert os.path.isfile(docx_path)
            assert os.path.getsize(docx_path) > 1000

            # Paste into PPTX
            pptx_path = os.path.join(tmp, "slides.pptx")
            result = writer.paste(screenshot, target=pptx_path)
            assert os.path.isfile(pptx_path)

            # Paste into Markdown
            md_path = os.path.join(tmp, "notes.md")
            result = writer.paste(screenshot, target=md_path, embed=True)
            assert os.path.isfile(md_path)
            with open(md_path) as f:
                assert "base64" in f.read()


# ---------------------------------------------------------------------------
# CLI tests
# ---------------------------------------------------------------------------

class TestCLI:
    def test_formats_command(self):
        from markitwrite.__main__ import main

        ret = main(["formats"])
        assert ret == 0

    def test_paste_command(self, tmp_path):
        from markitwrite.__main__ import main

        img_path = str(tmp_path / "img.png")
        with open(img_path, "wb") as f:
            f.write(_make_test_image())

        target = str(tmp_path / "out.docx")
        ret = main(["paste", img_path, "--to", target])
        assert ret == 0
        assert os.path.isfile(target)

    def test_no_command_returns_1(self):
        from markitwrite.__main__ import main

        ret = main([])
        assert ret == 1
