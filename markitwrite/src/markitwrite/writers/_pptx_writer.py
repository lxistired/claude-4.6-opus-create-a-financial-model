"""Writer for Microsoft PowerPoint (.pptx) presentations."""

from __future__ import annotations

import io
import sys
from typing import BinaryIO, Optional

from markitwrite._base_writer import DocumentWriter, WriteResult

_dependency_exc_info = None
try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    _dependency_exc_info = sys.exc_info()


class PptxWriter(DocumentWriter):
    """Insert images into PowerPoint (.pptx) presentations.

    Uses python-pptx under the hood.

    Position hints:
        {"slide": N} - Insert on slide N (1-indexed). 0 or missing = new slide.
        {"left": 1.0, "top": 1.0} - Position in inches from top-left.
        No position = centered on a new blank slide.

    Size hints:
        {"width": 8.0} - Width in inches (height auto-calculated).
        {"width": 6, "height": 4} - Both dimensions.
    """

    ACCEPTED_EXTENSIONS = [".pptx"]

    def insert_image(
        self,
        image_stream: BinaryIO,
        document_stream: Optional[BinaryIO] = None,
        target_format: str = ".pptx",
        position: Optional[dict] = None,
        size: Optional[dict] = None,
        **kwargs,
    ) -> WriteResult:
        if _dependency_exc_info is not None:
            raise ImportError(
                "python-pptx is required for PPTX writing. "
                "Install it with: pip install 'markitwrite[pptx]'"
            ) from _dependency_exc_info[1]

        # Load or create presentation
        if document_stream:
            document_stream.seek(0)
            prs = Presentation(document_stream)
        else:
            prs = Presentation()

        # Determine which slide to use
        slide = None
        if position and "slide" in position:
            slide_num = int(position["slide"])
            if 1 <= slide_num <= len(prs.slides):
                slide = prs.slides[slide_num - 1]

        if slide is None:
            # Add a new blank slide (layout index 6 is typically blank)
            blank_layout_idx = min(6, len(prs.slide_layouts) - 1)
            slide_layout = prs.slide_layouts[blank_layout_idx]
            slide = prs.slides.add_slide(slide_layout)

        # Determine position on slide
        left = Inches(float(position.get("left", 1.0))) if position else Inches(1.0)
        top = Inches(float(position.get("top", 1.0))) if position else Inches(1.0)

        # Determine sizing
        width = None
        height = None
        if size:
            if "width" in size:
                width = Inches(float(size["width"]))
            if "height" in size:
                height = Inches(float(size["height"]))
        if width is None and height is None:
            width = Inches(8.0)  # fill most of the slide width

        # Insert image
        image_stream.seek(0)
        slide.shapes.add_picture(
            image_stream, left, top, width=width, height=height
        )

        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        return WriteResult(
            output=output.read(),
            target_format=".pptx",
            images_inserted=1,
            metadata={
                "slide_count": len(prs.slides),
                "size": size or {"width": 8.0},
            },
        )
